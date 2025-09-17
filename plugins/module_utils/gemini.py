import google.generativeai as genai
import PIL.Image
from google.oauth2 import service_account
import os
import PyPDF2
import cv2
import requests
from pptx import Presentation

class Gemini:

    def __init__(self, service_account_info, project, location, model, prompt, context, src, multimodal, multimodal_type):
        self.service_account_info = service_account_info
        self.project = project
        self.location = location
        self.model = model
        self.prompt = prompt
        self.context = context
        self.src = src
        self.multimodal = multimodal
        self.multimodal_type = multimodal_type

    def leer_txt(self):
        with open(self.src, 'r') as f:
            texto = f.read()
            return texto
        
    def leer_pdf(self):
        with open(self.src, 'rb') as pdf_file_obj:
            pdf_reader = PyPDF2.PdfReader(pdf_file_obj)
            texto = ''
            for page in pdf_reader.pages:
                texto += page.extract_text()
        return texto

    def leer_ppt(self):
        prs = Presentation(self.src)
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                  texto += shape.text + "\n"
        return texto

    def descarga(self, ruta_local):
        response = requests.get(self.src, stream=True)
        if response.status_code == 200:
            with open(ruta_local, 'wb') as f:
                for chunk in response.iter_content(chunk_size=1024):
                    if chunk:
                        f.write(chunk)
            return ruta_local
        else:
            return None

    def leer_video(self):
        if self.src.startswith(("http://", "https://", "ftp://")):
            ruta_local = f"imagen{self.src.lower().endswith}"
            ruta_descargada = self.descarga(ruta_local)
            if ruta_descargada is None:
                return {"data": {"content": f"No se pudo descargar el video."}, 
                    "metadata": {"status": "failed"}}
            path = ruta_local
        else:
            path = self.src

        cap = cv2.VideoCapture(path)
        frames = []
        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break
            frames.append(frame)
        cap.release()
        num_frames = len(frames)
        step = max(1, num_frames // 10)
        frames_representativos = frames[::step]
        imagenes = []
        for frame in frames_representativos:
            imagen = PIL.Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
            imagenes.append(imagen)
        return imagenes

    def call_gemini(self):
        try:
            prompt_completo = f"{self.prompt} \n\n {self.context or ''}"
            credentials = service_account.Credentials.from_service_account_info(self.service_account_info)
            genai.configure(credentials=credentials)
            model = genai.GenerativeModel(model_name=self.model)

            if self.multimodal:
                
                if self.multimodal_type == 'imagen':
                    if not os.path.exists(self.src):
                        return {"data": {"content": f"El archivo '{self.src}' no fue encontrado."}, 
                                "metadata": {"status": "failed"}}
                    if not self.src.lower().endswith(('.png', '.jpg', '.jpeg')):
                        return {"data": {"content": f"Formato de imagen no soportado."}, 
                                "metadata": {"status": "failed"}}
                    if self.src.startswith(("http://", "https://", "ftp://")):
                      ruta_local = f"imagen{self.src.lower().endswith}"
                      ruta_descargada = self.descarga(ruta_local)
                      if ruta_descargada is None:
                        return {"data": {"content": f"No se pudo descargar la imagen."}, 
                                "metadata": {"status": "failed"}}
                      path = ruta_local
                    else:
                      path = self.src

                    image = PIL.Image.open(path)
                    response = model.generate_content([prompt_completo, image])
                    return {"data": {"content": response.text}, "metadata": {"status": "success"}}
                
                elif self.multimodal_type == 'archivo':
                    if not os.path.exists(self.src):
                        return {"data": {"content": f"El archivo '{self.src}' no fue encontrado."}, 
                                "metadata": {"status": "failed"}}
                    if self.src.endswith('.txt'):
                        texto = self.leer_txt()
                    elif self.src.endswith('.pdf'):
                        texto = self.leer_pdf()
                    elif self.src.endswith('.ppt') or self.src.endswith('.pptx'):
                        texto = self.leer_ppt()
                    else:
                        return {"data": {"content": f"Formato de archivo no soportado."}, 
                                "metadata": {"status": "failed"}}
                    response = model.generate_content([prompt_completo, texto])
                    return {"data": {"content": response.text}, "metadata": {"status": "success"}}
                             
                elif self.multimodal_type == 'video':
                    if not self.src.endswith('.mp4'):
                        return {"data": {"content": f"Formato de video no soportado."}, 
                                "metadata": {"status": "failed"}}
                    imagenes = self.leer_video()
                    if imagenes is None:
                        return {"data": {"content": f"No se pudo procesar el video."}, 
                                    "metadata": {"status": "failed"}}
                    respuesta = ""
                    for imagen in imagenes:
                        response = model.generate_content([prompt_completo, imagen])
                        respuesta += response.text + " "
                    return {"data": {"content": respuesta.strip()}, "metadata": {"status": "success"}}
                
            else:
                response = model.generate_content(prompt_completo)
                return {"data": {"content": response.text}, "metadata": {"status": "success"}}
            
        except Exception as e:
            return {
                "error": str(e),
                "data": None,
                "metadata": {"status": "failed"}
            }


